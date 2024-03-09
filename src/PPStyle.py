# Simple presentation from PowerPoint patterns
from pptx import Presentation
from pptx.util import Inches, Pt

from src import config


class PPPresentation:
    def __init__(
        self,
        text,
        theme,
        back_img_path=False,
        num_chapter=False,
    ) -> None:
        with open(text, encoding='utf-8') as f:
            self.text = f.readlines()

        self.root = Presentation(theme)
        self.num_chapter = 1 if num_chapter else False
        self.back_img_path = back_img_path

        self.layouts = {
            "title": 0,
            "chapter": 1,
        }  # take layouts from theme

    def background_img(self):
        title_slide_layout = self.root.slide_layouts[0]
        slide = self.root.slides.add_slide(title_slide_layout)
        left = top = Inches(0)
        pic = slide.shapes.add_picture(
            self.back_img_path,
            left,
            top,
            width=self.root.slide_width,
            height=self.root.slide_height,
        )

        # This moves it to the background
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

        return slide

    def add_table(self):
        pass

    def add_figure(self):
        pass

    def add_plot(self):
        pass

    def first_slide(self):
        # TODO: take structure from summarize text
        next_line_number = 0
        n_line = 0
        first_slide_lines = [False, False, False]
        while True:
            line = self.text[next_line_number]
            if (len(line) > 70) or (n_line > 3):
                break
            if len(line) != 0:
                first_slide_lines[n_line] = line.strip()
                n_line += 1
            next_line_number += 1

        title, subtitl, name = first_slide_lines
        title = config.NAME_PRESENTATION
        self.add_first_slide(title, subtitl, name)

        return next_line_number

    def add_first_slide(self, title, subtitle=False, name=False):
        layout = 0
        slide_layout = self.root.slide_layouts[layout]
        slide = self.root.slides.add_slide(slide_layout)

        title_, subtitle_, name_ = slide.placeholders

        title_.text = title
        if subtitle:
            subtitle_.text = subtitle
        if name:
            name_.text = name

    def chapter_slide(self, next_line_number):
        # TODO: take structure from summarize text
        n_line = 0
        chapter_slide_lines = [False, False, self.num_chapter]
        while True:
            line = self.text[next_line_number]
            if (len(line) > 100) or (n_line > 2):
                break
            if len(line) != 0:
                chapter_slide_lines[n_line] = line.strip()
                n_line += 1
            next_line_number += 1

        title, subtitl, num = chapter_slide_lines
        self.add_chapter_slide(title, subtitl, num)

        if self.num_chapter:
            self.num_chapter += 1

        return next_line_number

    def add_chapter_slide(self, title, subtitle=False, num=False):
        layout = 9
        slide_layout = self.root.slide_layouts[layout]
        slide = self.root.slides.add_slide(slide_layout)

        if num:
            title = " ".join([str(num), title])
        slide.shapes.title.text = title

        if subtitle:
            slide.placeholders[1].text = subtitle

    def text_slide(self, next_line_number):
        # TODO: take structure from summarize text
        line = self.text[next_line_number]

        title = line.split()[0]
        text = line

        self.add_text_slide(text, title)

    def add_text_slide(self, text, title=False):
        layout = 5 if title else 5
        slide_layout = self.root.slide_layouts[layout]
        slide = self.root.slides.add_slide(slide_layout)

        if title:
            slide.shapes.title.text = title

        slide.placeholders[1].text = text

    def create_presentation(self):
        next_line_number = self.first_slide()
        while next_line_number < len(self.text):
            line = self.text[next_line_number]
            if len(line.strip()) != 0:
                if len(line) < 60:
                    next_line_number = self.chapter_slide(next_line_number) - 1
                else:
                    self.text_slide(next_line_number)
            next_line_number += 1

        self.root.save(config.NAME_PRESENTATION)

